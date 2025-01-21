from pptx import Presentation
import os
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

def create_ppt_with_date_and_members(date, save_path, filename, members):
    script_dir = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(script_dir, "Presentation_template.pptx")
    
    prs = Presentation(template_path)
    
    slide = prs.slides[0]
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "5 minutes presentations"
    subtitle.text = f"Date: {date.strftime('%Y-%m-%d')}"
    
    for member in members:
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{member} - Previous week"
        
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{member} - This week"
    
    file_path = os.path.join(save_path, filename)

    if os.path.exists(file_path):
        os.remove(file_path)
    
    prs.save(file_path)
    return file_path

def upload_to_drive(service, file_path, filename, folder_id):
    file_metadata = {'name': filename, 'parents': [folder_id]}
    media = MediaFileUpload(file_path, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
    return file.get('id')

def create_shareable_link(service, file_id):
    permission = {'type': 'anyone', 'role': 'reader'}
    service.permissions().create(fileId=file_id, body=permission).execute()
    return f"https://drive.google.com/file/d/{file_id}/view"

if __name__ == "__main__":
    import argparse
    import json
    
    parser = argparse.ArgumentParser(description='Send a weekly reminder to Slack with a link to slides.')
    parser.add_argument('--date', required=True, help='Date')
    
    args = parser.parse_args()
    date = args.date

    # List of lab members
    lab_members = ["Guillaume", "Gautier", "Joanna", "Hannah", "Jaakko", "Ana", "Sujan", "Sarah", "Monika", "Marcela", "Iván", "Johanna"]

    # Get the values for the connection with Google Drive´s API
    SCOPES = ['https://www.googleapis.com/auth/drive.file']
    SERVICE_ACCOUNT_JSON = os.getenv('SERVICE_ACCOUNT_JSON')
    FOLDER_ID = os.getenv('FOLDER_ID')

    # # Get the credentials from Google Drive´s API
    credentials = service_account.Credentials.from_service_account_info(json.loads(SERVICE_ACCOUNT_JSON), scopes=SCOPES)
    service = build('drive', 'v3', credentials=credentials)

    # Folder to temporary store the presentaiton
    save_path = "presentations/"
    os.makedirs(save_path, exist_ok=True)

    # Dictionary to store dates and links
    links_dict = {}

    filename = f"{date}_5min_Presentation.pptx"
    file_path = create_ppt_with_date_and_members(date, save_path, filename, lab_members)
    file_id = upload_to_drive(service, file_path, filename, FOLDER_ID)
    shareable_link = create_shareable_link(service, file_id)
    
    print(shareable_link)