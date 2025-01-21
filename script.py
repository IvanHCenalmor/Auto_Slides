from pptx import Presentation
import os
import json
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

def create_ppt_with_date_and_members(date, save_path, filename, members):
    script_dir = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(script_dir, "Presentation_template.pptx")
    
    prs = Presentation(template_path)
    
    slide = prs.slides[0]
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "5 minutes presentations"
    subtitle.text = f"Date: {date.strftime('%Y-%m-%d')}"
    
    for member in members:
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{member} - Previous week"
        
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"{member} - This week"
    
    file_path = os.path.join(save_path, filename)

    if os.path.exists(file_path):
        os.remove(file_path)
    
    prs.save(file_path)
    return file_path

def upload_to_drive(service, file_path, filename, folder_id):
    file_metadata = {'name': filename, 'parents': [folder_id]}
    media = MediaFileUpload(file_path, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
    return file.get('id')

def create_shareable_link(service, file_id):
    permission = {'type': 'anyone', 'role': 'reader'}
    service.permissions().create(fileId=file_id, body=permission).execute()
    return f"https://drive.google.com/file/d/{file_id}/view"

if __name__ == "__main__":
    from datetime import datetime, timedelta

    # List of lab members
    lab_members = ["Guillaume", "Gautier", "Joanna", "Hannah", "Jaakko", "Ana", "Sujan", "Sarah", "Monika", "Marcela", "Iván", "Johanna"]

    SCOPES = ['https://www.googleapis.com/auth/drive.file']
    SERVICE_ACCOUNT_JSON = os.getenv('SERVICE_ACCOUNT_JSON')
    FOLDER_ID = os.getenv('FOLDER_ID')

    credentials = service_account.Credentials.from_service_account_info(json.loads(SERVICE_ACCOUNT_JSON), scopes=SCOPES)
    service = build('drive', 'v3', credentials=credentials)

    # Generate PowerPoint files with weekly dates and individual slides for each member
    start_date = datetime(2025, 1, 27)
    end_date = datetime(2025, 7, 1)
    current_date = start_date

    save_path = "presentations/"
    os.makedirs(save_path, exist_ok=True)

    # Dictionary to store dates and links
    links_dict = {}

    while current_date <= end_date:
        filename = f"{current_date.strftime('%Y-%m-%d')}_5min_Presentation.pptx"
        file_path = create_ppt_with_date_and_members(current_date, save_path, filename, lab_members)
        file_id = upload_to_drive(service, file_path, filename, FOLDER_ID)
        shareable_link = create_shareable_link(service, file_id)
        print(f"Presentation for {current_date.strftime('%Y-%m-%d')} uploaded: {shareable_link}")
        
        # Store the date and link in the dictionary
        links_dict[current_date.strftime('%Y-%m-%d')] = shareable_link
        current_date += timedelta(weeks=1)

    # Write the dictionary to a JSON file
    with open('presentation_links.json', 'w') as json_file:
        json.dump(links_dict, json_file, indent=4)